"""
build.py — Generate thesis seminar presentation using official TAU visual identity
Brand color: #4E008E (Tampere University purple)

Usage:
    /tmp/pptx-venv/bin/python build.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree
import copy, os

# ---------------------------------------------------------------------------
# Official TAU brand colors
# ---------------------------------------------------------------------------
TAU_PURPLE  = RGBColor(0x4E, 0x00, 0x8E)   # #4E008E — official brand color
TAU_LIGHT   = RGBColor(0x9B, 0x59, 0xC7)   # lighter purple for accents
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY   = RGBColor(0x1A, 0x1A, 0x1A)
MID_GREY    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY  = RGBColor(0xF5, 0xF0, 0xFA)   # very light purple tint
OFF_WHITE   = RGBColor(0xFA, 0xF8, 0xFC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

LOGO_WHITE  = os.path.join(os.path.dirname(__file__), "tau_logo_white.png")
LOGO_PURPLE = os.path.join(os.path.dirname(__file__), "tau_logo_purple.png")

# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, left, top, width, height, fill=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=16, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_logo(slide, white=True, left=10.8, top=0.18, width=2.2):
    path = LOGO_WHITE if white else LOGO_PURPLE
    # logo aspect: 266.32 × 84.8 ≈ 3.14:1
    height = width / 3.14
    slide.shapes.add_picture(path, Inches(left), Inches(top),
                              Inches(width), Inches(height))


def header_bar(slide, title, subtitle=None, bar_h=1.1):
    """Purple header bar with white title."""
    rect(slide, 0, 0, 13.33, bar_h, fill=TAU_PURPLE)
    add_text(slide, title,
             left=0.4, top=0.12, width=10.0, height=0.7,
             size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 left=0.4, top=0.72, width=10.5, height=0.38,
                 size=12, color=RGBColor(0xDD, 0xBB, 0xFF))
    add_logo(slide, white=True, left=10.85, top=0.17, width=2.1)


def footer_bar(slide, page, total=9):
    rect(slide, 0, 7.22, 13.33, 0.28, fill=TAU_PURPLE)
    add_text(slide, "Tampere University  |  Nhan Chau  |  April 2026",
             left=0.3, top=7.24, width=11.0, height=0.24,
             size=8, color=RGBColor(0xDD, 0xBB, 0xFF), align=PP_ALIGN.LEFT)
    add_text(slide, f"{page} / {total}",
             left=12.5, top=7.24, width=0.7, height=0.24,
             size=8, color=WHITE, align=PP_ALIGN.RIGHT)


def bullet_list(slide, items, left, top, width, spacing=0.42,
                size=15, color=DARK_GREY, indent_color=None, marker="•"):
    y = top
    for item in items:
        full = f"{marker}  {item}"
        add_text(slide, full, left=left, top=y,
                 width=width, height=spacing * 0.92,
                 size=size, color=color)
        y += spacing
    return y


def draw_table(slide, headers, rows, left, top, width, col_widths,
               row_h=0.38, hdr_size=12, cell_size=11):
    # header
    rect(slide, left, top, width, row_h, fill=TAU_PURPLE)
    x = left
    for i, h in enumerate(headers):
        add_text(slide, h, left=x + 0.06, top=top + 0.05,
                 width=col_widths[i] - 0.1, height=row_h - 0.08,
                 size=hdr_size, bold=True, color=WHITE)
        x += col_widths[i]
    # rows
    for r, row in enumerate(rows):
        ry = top + row_h * (r + 1)
        bg = LIGHT_GREY if r % 2 == 0 else OFF_WHITE
        rect(slide, left, ry, width, row_h, fill=bg)
        x = left
        for c, cell in enumerate(row):
            add_text(slide, cell, left=x + 0.06, top=ry + 0.04,
                     width=col_widths[c] - 0.1, height=row_h - 0.06,
                     size=cell_size, color=DARK_GREY)
            x += col_widths[c]


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def s01_title(prs):
    slide = blank_slide(prs)

    # Full purple background
    rect(slide, 0, 0, 13.33, 7.5, fill=TAU_PURPLE)

    # Large "T" decorative element (subtle)
    rect(slide, 10.5, 0, 2.83, 7.5, fill=RGBColor(0x5A, 0x00, 0xA8))

    # White logo top-left
    add_logo(slide, white=True, left=0.45, top=0.35, width=3.0)

    # Title
    add_text(slide,
             "Design and Implementation of a\nLangGraph-Based Marketing Analytics Agent",
             left=0.45, top=1.6, width=9.8, height=2.2,
             size=34, bold=True, color=WHITE)

    # Subtitle
    add_text(slide,
             "An Exploration of Agentic Patterns for Autonomous Data Analysis",
             left=0.45, top=3.75, width=9.8, height=0.6,
             size=16, color=RGBColor(0xDD, 0xBB, 0xFF))

    # Divider
    rect(slide, 0.45, 4.55, 4.0, 0.04, fill=RGBColor(0x9B, 0x59, 0xC7))

    # Author block
    add_text(slide, "Nhan Chau",
             left=0.45, top=4.75, width=6.0, height=0.5,
             size=20, bold=True, color=WHITE)
    add_text(slide,
             "Bachelor's Thesis  •  Tampere University  •  April 2026\n"
             "Examiner: Professor Ahmed Farooq",
             left=0.45, top=5.25, width=9.5, height=0.8,
             size=13, color=RGBColor(0xDD, 0xBB, 0xFF))

    # Page
    add_text(slide, "1 / 9", left=12.55, top=7.18, width=0.7, height=0.25,
             size=8, color=RGBColor(0xDD, 0xBB, 0xFF), align=PP_ALIGN.RIGHT)


def s02_intro(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Introduction")
    footer_bar(slide, 2, total=11)

    sentences = [
        ("What",
         "This thesis designs a conversational AI agent that lets marketing analysts "
         "query data and produce reports in natural language — without writing SQL or using BI tools."),
        ("How",
         "The system is built on LangGraph using the ReAct pattern: a large language model "
         "drives a loop of reasoning and tool calls against a live marketing database."),
        ("Scope",
         "The work covers system architecture, three key design decisions, and a planned "
         "evaluation across six representative scenarios."),
    ]

    y = 1.6
    for label, body in sentences:
        rect(slide, 0.4, y, 1.1, 0.95, fill=TAU_PURPLE)
        add_text(slide, label, left=0.42, top=y + 0.28,
                 width=1.06, height=0.4,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, body, left=1.7, top=y + 0.15,
                 width=11.1, height=0.72, size=15, color=DARK_GREY)
        y += 1.55


def s02_agenda(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Agenda")
    footer_bar(slide, 3, total=11)

    items = [
        "Introduction",
        "Problem & Motivation",
        "Research Questions",
        "System Architecture",
        "Key Design #1 — Smart Agent, Simple Tools",
        "Key Design #2 — Skills System",
        "Key Design #3 — Parallel Multi-Agent Orchestration",
        "Evaluation Plan",
        "Conclusions",
    ]
    bullet_list(slide, items, left=1.0, top=1.35, width=11.3,
                spacing=0.54, size=17)


def s03_problem(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Problem & Motivation",
               "Marketing analytics today is slow, manual, and fragmented")
    footer_bar(slide, 4, total=11)

    pain = [
        "Analyst receives a question → writes SQL → cleans data → writes report",
        "Each cycle takes hours; repeated for every ad-hoc request",
        "BI tools help with standard dashboards — not free-form investigation",
    ]
    bullet_list(slide, pain, left=0.6, top=1.28, width=12.1,
                spacing=0.55, size=16)

    # Hook box
    rect(slide, 0.5, 3.2, 12.3, 1.1, fill=RGBColor(0x3A, 0x00, 0x6A))
    add_text(slide,
             "Modern LLMs can write SQL, reason about results, and generate structured reports.\n"
             "Can a conversational agent replace the full manual analytics cycle?",
             left=0.75, top=3.32, width=11.8, height=0.88,
             size=15, bold=False, color=WHITE)

    add_text(slide, "Thesis goal",
             left=0.5, top=4.6, width=2.5, height=0.35,
             size=13, bold=True, color=TAU_PURPLE)
    add_text(slide,
             "Design a LangGraph-based conversational AI agent that allows marketing analysts "
             "to query data and produce reports using natural language — without SQL or BI tools.",
             left=0.5, top=4.93, width=12.3, height=0.7,
             size=13, color=DARK_GREY)


def s04_rqs(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Research Questions")
    footer_bar(slide, 5, total=11)

    rqs = [
        ("RQ1",
         "How can a ReAct-based LLM agent be designed to reliably perform\n"
         "multi-step marketing data analysis?"),
        ("RQ2",
         "How can parallel multi-agent orchestration reduce analysis latency\n"
         "and support cross-platform comparisons?"),
        ("RQ3",
         "How do structured prompt templates (skills) improve the reliability\n"
         "and consistency of an LLM agent for domain-specific workflows?"),
    ]

    y = 1.3
    for label, text in rqs:
        rect(slide, 0.4, y, 1.1, 1.0, fill=TAU_PURPLE)
        add_text(slide, label, left=0.42, top=y + 0.25,
                 width=1.06, height=0.5,
                 size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rect(slide, 1.5, y, 11.3, 1.0, fill=LIGHT_GREY)
        add_text(slide, text, left=1.7, top=y + 0.15,
                 width=11.0, height=0.82, size=15, color=DARK_GREY)
        y += 1.5


def s05_architecture(prs):
    slide = blank_slide(prs)
    header_bar(slide, "System Architecture",
               "Four layers from user message to final answer")
    footer_bar(slide, 6, total=11)

    headers = ["Layer", "Technology", "Responsibility"]
    rows = [
        ["API layer",   "FastAPI",      "Accepts user messages; returns final response"],
        ["Agent layer", "LangGraph",    "Runs the ReAct loop; manages conversation state"],
        ["Tool layer",  "Python / SQL", "3 tools: schema discovery & SQL execution"],
        ["Data layer",  "PostgreSQL",   "Marketing database — read-only access"],
    ]
    draw_table(slide, headers, rows,
               left=0.4, top=1.28, width=12.5,
               col_widths=[1.8, 2.0, 8.7], row_h=0.42)

    # Flow diagram
    flow = [("User", TAU_PURPLE), ("FastAPI", TAU_LIGHT),
            ("LangGraph\nReAct", TAU_PURPLE), ("Tools", TAU_LIGHT),
            ("PostgreSQL", TAU_PURPLE)]
    fx = 0.55
    fy = 4.1
    bw, bh = 2.1, 0.82
    gap = 0.28
    for i, (label, col) in enumerate(flow):
        rect(slide, fx, fy, bw, bh, fill=col)
        lines = label.split("\n")
        if len(lines) == 2:
            add_text(slide, lines[0], left=fx, top=fy + 0.1,
                     width=bw, height=0.3, size=12, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
            add_text(slide, lines[1], left=fx, top=fy + 0.42,
                     width=bw, height=0.3, size=12, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, label, left=fx, top=fy + 0.24,
                     width=bw, height=0.38, size=12, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(slide, "→", left=fx + bw + 0.02, top=fy + 0.2,
                     width=gap, height=0.4, size=16, bold=True,
                     color=TAU_PURPLE, align=PP_ALIGN.CENTER)
        fx += bw + gap + 0.02

    add_text(slide,
             "The lead agent node calls the LLM; the router directs to the tool executor "
             "if the model emitted a tool call; the loop continues until a final text response.",
             left=0.4, top=5.18, width=12.5, height=0.6,
             size=12, italic=True, color=MID_GREY)


def s06_tools(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Key Design #1 — Smart Agent, Simple Tools",
               "Only add a tool when it crosses a boundary the model cannot cross alone")
    footer_bar(slide, 7, total=11)

    headers = ["Tool", "Category", "Purpose"]
    rows = [
        ["`list_tables`",    "Schema", "List all tables with row counts"],
        ["`describe_table`", "Schema", "Show schema, sample rows, value distributions"],
        ["`query_data`",     "Data",   "Execute read-only SQL; return results as Markdown"],
    ]
    draw_table(slide, headers, rows,
               left=0.4, top=1.28, width=12.5,
               col_widths=[2.0, 1.6, 8.9], row_h=0.46)

    # Two insight boxes
    for i, (title, body) in enumerate([
        ("Fewer tools = higher reliability",
         "Fewer tools means fewer decision points for the model — less latency, fewer wrong choices."),
        ("Description is prompt engineering",
         "Changing only the tool description (not the code) changes how the model behaves."),
    ]):
        bx = 0.4 + i * 6.3
        rect(slide, bx, 3.65, 6.0, 1.45, fill=LIGHT_GREY)
        rect(slide, bx, 3.65, 6.0, 0.36, fill=TAU_PURPLE)
        add_text(slide, title, left=bx + 0.12, top=3.68,
                 width=5.76, height=0.3, size=12, bold=True, color=WHITE)
        add_text(slide, body, left=bx + 0.12, top=4.07,
                 width=5.76, height=0.95, size=13, color=DARK_GREY)

    add_text(slide,
             "The model cannot query the database on its own — hence the three tools. "
             "All other analytical work (reasoning, formatting, summarising) stays in the model.",
             left=0.4, top=5.3, width=12.5, height=0.55,
             size=12, italic=True, color=MID_GREY)


def s07_skills(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Key Design #2 — Skills System",
               "Structured workflow recipes injected into the agent's session context at runtime")
    footer_bar(slide, 8, total=11)

    # Left: problem + solution
    add_text(slide, "Problem", left=0.4, top=1.28, width=2.0, height=0.35,
             size=14, bold=True, color=TAU_PURPLE)
    add_text(slide,
             "The ReAct loop is flexible but unpredictable for complex multi-step workflows. "
             "Without guidance, the model may skip steps or produce inconsistent output.",
             left=0.4, top=1.62, width=6.1, height=0.8, size=13, color=DARK_GREY)

    add_text(slide, "Solution — a Skill", left=0.4, top=2.6, width=3.5, height=0.35,
             size=14, bold=True, color=TAU_PURPLE)
    items = [
        "Trigger keyword  (e.g. /weekly-report)",
        "Numbered step sequence to follow exactly",
        "Output format: section headings, table structure",
        "Injected into session context — no new tools needed",
    ]
    bullet_list(slide, items, left=0.5, top=2.95, width=6.0,
                spacing=0.44, size=13)

    # Right: skill file example box
    rect(slide, 6.85, 1.28, 6.1, 4.2, fill=RGBColor(0x1A, 0x00, 0x30))
    add_text(slide, "weekly-report skill  (excerpt)",
             left=6.95, top=1.32, width=5.9, height=0.32,
             size=10, bold=True, color=RGBColor(0xDD, 0xBB, 0xFF))
    code = (
        "trigger: /weekly-report\n"
        "tables:  daily_metrics, campaigns\n"
        "---\n"
        "1. Query last 7 days grouped by platform\n"
        "2. Compute week-over-week changes:\n"
        "   spend, CTR, CPC, ROAS\n"
        "3. Identify top 5 / bottom 5 campaigns\n"
        "4. Output: executive summary, tables,\n"
        "   trend analysis"
    )
    add_text(slide, code, left=6.95, top=1.68, width=5.9, height=3.6,
             size=12, color=RGBColor(0xCC, 0xFF, 0xCC))

    # Analogy
    rect(slide, 0.4, 5.5, 12.5, 0.65, fill=RGBColor(0x3A, 0x00, 0x6A))
    add_text(slide,
             "Analogy: the agent already knows how to cook — the skill gives it a recipe "
             "so every weekly report is complete and consistent.",
             left=0.6, top=5.58, width=12.2, height=0.5,
             size=13, color=WHITE)


def s08_multiagent(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Key Design #3 — Parallel Multi-Agent Orchestration",
               "Lead agent spawns independent subagents; results synthesised from text responses")
    footer_bar(slide, 9, total=11)

    # Left column
    add_text(slide, "How it works", left=0.4, top=1.28, width=5.5, height=0.35,
             size=14, bold=True, color=TAU_PURPLE)
    steps = [
        "Lead agent calls research_agent with N tasks",
        "Each task runs as an independent LangGraph graph",
        "Coroutines launched with asyncio.gather() in parallel",
        "Subagents return structured text results",
        "Lead agent synthesises into unified report",
    ]
    bullet_list(slide, steps, left=0.4, top=1.65, width=6.1,
                spacing=0.48, size=13)

    # Right: result box
    rect(slide, 6.9, 1.28, 6.05, 3.3, fill=LIGHT_GREY)
    rect(slide, 6.9, 1.28, 6.05, 0.4, fill=TAU_PURPLE)
    add_text(slide, "Expected latency benefit",
             left=7.0, top=1.31, width=5.85, height=0.34,
             size=12, bold=True, color=WHITE)
    add_text(slide,
             "Cross-platform comparison (Google vs Meta Q1)\n\n"
             "Sequential execution:    O(n) — platforms run one by one\n"
             "Parallel execution:        O(1) — all platforms run at once\n\n"
             "Benefit grows with the number of platforms compared.",
             left=7.0, top=1.76, width=5.85, height=2.7,
             size=13, color=DARK_GREY)

    # Design note
    rect(slide, 0.4, 4.8, 12.5, 0.75, fill=RGBColor(0x3A, 0x00, 0x6A))
    add_text(slide,
             "Design rule:  use subagents when each task would require more than 3 tool calls on its own — "
             "that's when parallelism overhead is outweighed by the latency saving.",
             left=0.6, top=4.88, width=12.2, height=0.6,
             size=12, color=WHITE)


def s09_evaluation(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Evaluation Plan",
               "Six scenarios covering primary use cases; qualitative correctness + description experiment")
    footer_bar(slide, 10, total=11)

    headers = ["#", "Scenario", "Type", "What it validates"]
    rows = [
        ["1", "Simple question: spend last month",       "Single query",  "Schema discovery + SQL generation"],
        ["2", "Follow-up: break down by campaign type",  "Follow-up",     "Context reuse across turns"],
        ["3", "Trend analysis: 30-day platform spend",   "Multi-query",   "Date-range SQL + written summary"],
        ["4", "Investigation: why did ROAS drop?",       "Multi-step",    "Autonomous multi-step reasoning"],
        ["5", "/weekly-report skill",                    "Skill-driven",  "Skill vs. no-skill consistency"],
        ["6", "Compare Google vs Meta Q1",               "Multi-agent",   "Parallel subagent latency benefit"],
    ]
    draw_table(slide, headers, rows,
               left=0.4, top=1.28, width=12.5,
               col_widths=[0.35, 3.3, 1.7, 7.15],
               row_h=0.42, hdr_size=11, cell_size=10)

    # Prompt experiment
    rect(slide, 0.4, 5.38, 12.5, 0.75, fill=RGBColor(0x3A, 0x00, 0x6A))
    add_text(slide, "Key experiment:",
             left=0.6, top=5.44, width=2.2, height=0.3,
             size=12, bold=True, color=RGBColor(0xDD, 0xBB, 0xFF))
    add_text(slide,
             "Same tool code, two descriptions (minimal vs. full). "
             "Hypothesis: full description eliminates schema hallucination and prevents write-operation attempts.",
             left=2.7, top=5.44, width=10.1, height=0.6,
             size=12, color=WHITE)


def s10_conclusions(prs):       # renamed from s10 — no next steps
    slide = blank_slide(prs)
    header_bar(slide, "Conclusions")
    footer_bar(slide, 11, total=11)

    # Left: contributions
    add_text(slide, "Design contributions", left=0.4, top=1.28,
             width=5.9, height=0.38, size=15, bold=True, color=TAU_PURPLE)
    contribs = [
        "Proposed LangGraph ReAct architecture with 3 minimal tools",
        "Parallel multi-agent design via asyncio.gather()",
        "Skills system — reliability without fine-tuning",
        "Tool description design principle (description > code)",
    ]
    bullet_list(slide, contribs, left=0.4, top=1.68, width=5.9,
                spacing=0.5, size=13)

    # Right: RQ answers
    rect(slide, 6.7, 1.28, 6.25, 4.4, fill=LIGHT_GREY)
    rect(slide, 6.7, 1.28, 6.25, 0.38, fill=TAU_PURPLE)
    add_text(slide, "Research question answers",
             left=6.82, top=1.31, width=6.05, height=0.32,
             size=12, bold=True, color=WHITE)

    rqa = [
        ("RQ1", "System prompt + skills + tool error messages\n→ reliable multi-step analysis"),
        ("RQ2", "asyncio.gather() parallel subagents\n→ O(1) latency vs O(n) sequential"),
        ("RQ3", "Skills provide explicit step sequence\n→ zero ambiguity for the model"),
    ]
    y = 1.75
    for label, text in rqa:
        add_text(slide, label, left=6.82, top=y, width=0.9, height=0.7,
                 size=12, bold=True, color=TAU_PURPLE)
        add_text(slide, text, left=7.72, top=y, width=5.1, height=0.7,
                 size=12, color=DARK_GREY)
        y += 1.1

    # Thank you
    rect(slide, 0, 6.48, 13.33, 1.02, fill=TAU_PURPLE)
    add_text(slide, "Thank you  —  Questions welcome",
             left=0, top=6.65, width=13.33, height=0.5,
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = new_prs()
    s01_title(prs)
    s02_intro(prs)
    s02_agenda(prs)
    s03_problem(prs)
    s04_rqs(prs)
    s05_architecture(prs)
    s06_tools(prs)
    s07_skills(prs)
    s08_multiagent(prs)
    s09_evaluation(prs)
    s10_conclusions(prs)

    out = "thesis-presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides._sldIdLst.__len__()} slides)")


if __name__ == "__main__":
    build()
